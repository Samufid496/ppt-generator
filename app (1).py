import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from copy import deepcopy
import io
import os
from PIL import Image

st.set_page_config(page_title="PPT Generator | Pan Gulf Marketing", page_icon="📊", layout="centered")

st.markdown("""
<style>
    .stApp { background-color: #F4F6FB; }
    .block-container { padding-top: 2rem; max-width: 780px; }
    .banner { background: linear-gradient(135deg, #00ADEE, #1E2761); border-radius: 12px; padding: 28px 32px; margin-bottom: 28px; color: white; }
    .banner h1 { color: white; margin: 0; font-size: 2rem; }
    .banner p  { color: rgba(255,255,255,0.85); margin: 4px 0 0; font-size: 1rem; }
    .step-card { background: white; border-radius: 10px; padding: 20px 24px; margin-bottom: 16px; border-left: 4px solid #00ADEE; box-shadow: 0 2px 8px rgba(0,0,0,0.06); }
    .step-label { font-size: 0.75rem; font-weight: 700; color: #00ADEE; text-transform: uppercase; letter-spacing: 1px; margin-bottom: 4px; }
    .tip { background: #EBF8FF; border-radius: 8px; padding: 10px 14px; font-size: 0.85rem; color: #1E2761; margin-top: 8px; }
</style>
""", unsafe_allow_html=True)

st.markdown("""
<div class="banner">
    <h1>📊 PPT Generator</h1>
    <p>Pan Gulf Marketing — Discounter Slides | Upload Excel → Get PowerPoint</p>
</div>
""", unsafe_allow_html=True)

def fmt_num(val):
    try:
        f = float(val)
        return str(int(f)) if f == int(f) else str(round(f, 4))
    except:
        return str(val)

def set_cell_text(cell, text, font_size=None):
    tf = cell.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = str(text)
            if font_size:
                run.font.size = Pt(font_size)
            return
        run = para.add_run()
        run.text = str(text)
        if font_size:
            run.font.size = Pt(font_size)
        return

def get_placeholder_pos(slide):
    """Find imgProduct shape, grab its position/size, remove it from slide."""
    for shape in slide.shapes:
        if shape.name == 'imgProduct':
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            shape._element.getparent().remove(shape._element)
            return left, top, width, height
    # Fallback position if shape not found
    return Inches(3.4), Inches(1.53), Inches(6.54), Inches(3.96)

def update_slide(slide, row, image_bytes=None):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == 'txtDescription':
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = str(row['Material Description'])
        if shape.has_table:
            table = shape.table
            dr = table.rows[1]
            vals = [
                str(row['Material']),
                str(row['Sales Unit']),
                str(int(float(row['Barcode']))),
                fmt_num(row['R.P']),
                fmt_num(row['Net 25%']),
                fmt_num(row['Add 5% on Net']),
            ]
            for i, cell in enumerate(dr.cells):
                set_cell_text(cell, vals[i], 11 if i == 2 else None)

    left, top, width, height = get_placeholder_pos(slide)
    if image_bytes:
        slide.shapes.add_picture(image_bytes, left, top, width, height)

def duplicate_slide(prs):
    template = prs.slides[0]
    new_slide = prs.slides.add_slide(template.slide_layout)
    t_spTree = template._element.find(qn('p:cSld')).find(qn('p:spTree'))
    n_spTree = new_slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    for child in list(n_spTree):
        n_spTree.remove(child)
    for child in list(t_spTree):
        n_spTree.append(deepcopy(child))
    return new_slide

def generate(template_bytes, df, images_dict):
    prs = Presentation(io.BytesIO(template_bytes))
    total = len(df)
    mat0 = str(df.iloc[0]['Material'])
    img0 = io.BytesIO(images_dict[mat0]) if mat0 in images_dict else None
    update_slide(prs.slides[0], df.iloc[0], img0)
    yield 1, total
    for i in range(1, total):
        row = df.iloc[i]
        new_slide = duplicate_slide(prs)
        mat = str(row['Material'])
        img = io.BytesIO(images_dict[mat]) if mat in images_dict else None
        update_slide(new_slide, row, img)
        yield i + 1, total
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    yield output, total

# ── Step 1: Template ──────────────────────────────────────────────────────────
st.markdown('<div class="step-card"><div class="step-label">Step 1</div>', unsafe_allow_html=True)
st.markdown("**Upload your Template PPT** (master slide)")
template_file = st.file_uploader("", type=["pptx"], key="template", label_visibility="collapsed")
if template_file:
    st.success(f"✅ Template loaded: `{template_file.name}`")
st.markdown('<div class="tip">💡 Use <b>Discounter_working.pptx</b> — it already has the <code>imgProduct</code> container built in.</div>', unsafe_allow_html=True)
st.markdown('</div>', unsafe_allow_html=True)

# ── Step 2: Excel ─────────────────────────────────────────────────────────────
st.markdown('<div class="step-card"><div class="step-label">Step 2</div>', unsafe_allow_html=True)
st.markdown("**Upload your Excel file** (product data)")
excel_file = st.file_uploader("", type=["xlsx", "xls"], key="excel", label_visibility="collapsed")
df_preview = None
if excel_file:
    try:
        df_preview = pd.read_excel(excel_file)
        st.success(f"✅ Excel loaded: **{len(df_preview)} rows** found")
        with st.expander("Preview data (first 5 rows)"):
            st.dataframe(df_preview.head())
    except Exception as e:
        st.error(f"Could not read Excel: {e}")
st.markdown('<div class="tip">💡 Must have columns: <code>Material Description, Material, Sales Unit, Barcode, R.P, Net 25%, Add 5% on Net</code></div>', unsafe_allow_html=True)
st.markdown('</div>', unsafe_allow_html=True)

# ── Step 3: Images ────────────────────────────────────────────────────────────
st.markdown('<div class="step-card"><div class="step-label">Step 3 — Optional</div>', unsafe_allow_html=True)
st.markdown("**Upload product images** (optional)")
image_files = st.file_uploader("", type=["jpg","jpeg","png"], accept_multiple_files=True, key="images", label_visibility="collapsed")
images_dict = {}
if image_files:
    for f in image_files:
        name = os.path.splitext(f.name)[0]
        images_dict[name] = f.read()
    st.success(f"✅ {len(images_dict)} image(s) uploaded")
    with st.expander("See uploaded images"):
        cols = st.columns(4)
        for idx, (name, data) in enumerate(images_dict.items()):
            with cols[idx % 4]:
                st.image(data, caption=name, use_column_width=True)
st.markdown('<div class="tip">💡 Name each image exactly like the Material code — e.g. <code>MI-PIN200-06CY.jpg</code><br>Images will be placed inside the <b>imgProduct container</b> in your template.</div>', unsafe_allow_html=True)
st.markdown('</div>', unsafe_allow_html=True)

# ── Step 4: Generate ──────────────────────────────────────────────────────────
st.markdown('<div class="step-card"><div class="step-label">Step 4</div>', unsafe_allow_html=True)
st.markdown("**Generate your presentation**")
output_name = st.text_input("Output filename", value="Output.pptx")
if not output_name.endswith(".pptx"):
    output_name += ".pptx"

generate_btn = st.button("⚡  Generate Presentation", type="primary",
    disabled=not (template_file and df_preview is not None), use_container_width=True)

if generate_btn:
    required_cols = ['Material Description','Material','Sales Unit','Barcode','R.P','Net 25%','Add 5% on Net']
    missing = [c for c in required_cols if c not in df_preview.columns]
    if missing:
        st.error(f"❌ Excel is missing columns: {', '.join(missing)}")
    else:
        template_bytes = template_file.getvalue()
        progress_bar = st.progress(0, text="Starting…")
        status = st.empty()
        output_bytes = None
        total_slides = 0
        gen = generate(template_bytes, df_preview, images_dict)
        for result in gen:
            if isinstance(result[0], int):
                done, total = result
                progress_bar.progress(int(done/total*100), text=f"Processing slide {done} of {total}…")
            else:
                output_bytes, total_slides = result
        progress_bar.progress(100, text="✅ Done!")
        status.success(f"🎉 {total_slides} slides generated successfully!")
        st.download_button(label="📥  Download Presentation", data=output_bytes,
            file_name=output_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True, type="primary")

st.markdown('</div>', unsafe_allow_html=True)
st.markdown("---")
st.markdown("<p style='text-align:center; color:#999; font-size:0.8rem;'>Pan Gulf Marketing © 2025 · Built with Streamlit</p>", unsafe_allow_html=True)
